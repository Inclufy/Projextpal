"""
PPTX export for the PRINCE2 Highlight Report.

Generates a deck that mirrors Yanmar's "YEU Projects - Highlight Report Template.pptx":

    Slide 1: Cover  -- "<Project> Highlight Report  /  Last update: <date>"
    Slide 2: Body
        - Header block: Report to / Sponsor / Project Manager / Senior Supplier / Objective(s)
        - 4-axis RAG (Budget / Planning / Resources / Overall)
        - Monthly timeline bar (rolling 10 months, phase rows)
        - Financials table (Costs / Internal / External / Contingency
                            vs Budget / Actuals / ETC / Variance)
        - 5x5 Risk Map (probability x impact)
        - Highlights / Lowlights narrative
        - Issues / Risks lists

The implementation is deliberately self-contained: it builds the deck from
scratch with python-pptx so we don't depend on a binary template file in the
repo. Yanmar's template uses a plain white background, dark text, and 3-state
RAG dots -- which we replicate with shape primitives.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import date, timedelta
from io import BytesIO
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# RGB palette matching Yanmar templates / common steerco vocabulary.
RAG_COLORS = {
    "green": RGBColor(0x2E, 0x7D, 0x32),
    "amber": RGBColor(0xF9, 0xA8, 0x25),
    "red":   RGBColor(0xC6, 0x28, 0x28),
}
GREY_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREY_BORDER = RGBColor(0xBD, 0xBD, 0xBD)
TEXT_DARK = RGBColor(0x21, 0x21, 0x21)
TEXT_MUTED = RGBColor(0x61, 0x61, 0x61)
TIMELINE_FILL = RGBColor(0x1E, 0x88, 0xE5)


@dataclass
class HighlightReportData:
    """Plain DTO so the renderer is decoupled from Django ORM."""

    project_name: str = ""
    report_date: date | None = None
    period_label: str = ""
    report_month: int | None = None
    report_to: str = ""
    sponsor: str = ""
    project_manager: str = ""
    senior_supplier: str = ""
    objectives: str = ""

    rag_overall: str = "green"
    rag_budget: str = "green"
    rag_planning: str = "green"
    rag_resources: str = "green"

    # Each phase: {"name": "Prepare", "start": date, "end": date}
    phases: list[dict] = field(default_factory=list)

    # Each finance row: {"label": "Costs", "budget": Decimal, "actuals": Decimal,
    #                    "etc": Decimal, "variance": Decimal}
    financial_rows: list[dict] = field(default_factory=list)

    # Each risk: {"id": "R1", "probability": 1-5, "impact": 1-5, "name": str}
    risks_map: list[dict] = field(default_factory=list)

    highlights: str = ""
    lowlights: str = ""
    issues: list[str] = field(default_factory=list)
    risks: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_text(slide, x, y, w, h, text, *, bold=False, size=10,
              color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_rag_dot(slide, x, y, status, *, size=Inches(0.25)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RAG_COLORS.get(status, RAG_COLORS["green"])
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.width = Pt(1.5)
    return shape


def _add_panel(slide, x, y, w, h, *, fill=GREY_BG, border=GREY_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _months_window(today: date, n: int = 10) -> list[date]:
    """Return the first day of the n months starting 2 months before `today`."""
    start = date(today.year, today.month, 1) - timedelta(days=60)
    start = date(start.year, start.month, 1)
    months = []
    y, m = start.year, start.month
    for _ in range(n):
        months.append(date(y, m, 1))
        m += 1
        if m > 12:
            m = 1
            y += 1
    return months


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_cover(prs: Presentation, data: HighlightReportData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text(
        slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.8),
        f"{data.project_name or 'Project'}",
        bold=True, size=36, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.6),
        "Highlight Report",
        size=24, align=PP_ALIGN.CENTER, color=TEXT_MUTED,
    )
    update = data.report_date.strftime("%-d %B %Y") if data.report_date else ""
    _add_text(
        slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
        f"Last update: {update}" if update else "",
        size=14, align=PP_ALIGN.CENTER, color=TEXT_MUTED,
    )


def _build_header_block(slide, data: HighlightReportData) -> None:
    # Title bar
    _add_panel(slide, Inches(0.3), Inches(0.25), Inches(12.7), Inches(0.55),
               fill=RGBColor(0x1A, 0x1A, 0x1A), border=RGBColor(0x1A, 0x1A, 0x1A))
    _add_text(
        slide, Inches(0.45), Inches(0.30), Inches(12.5), Inches(0.45),
        f"Highlight Report — {data.project_name or 'Project'}",
        bold=True, size=18, color=RGBColor(0xFF, 0xFF, 0xFF),
    )

    # Left meta block (Report to / Sponsor / etc.)
    rows = [
        ("Report to", data.report_to),
        ("Sponsor", data.sponsor),
        ("Project Manager", data.project_manager),
        ("Senior Supplier", data.senior_supplier),
        ("Objective(s)", data.objectives),
    ]
    y = Inches(0.95)
    for label, value in rows:
        _add_text(slide, Inches(0.35), y, Inches(1.5), Inches(0.28),
                  f"{label}:", bold=True, size=10)
        _add_text(slide, Inches(1.85), y, Inches(4.8), Inches(0.28),
                  value or "—", size=10)
        y += Inches(0.30)


def _build_rag_block(slide, data: HighlightReportData) -> None:
    """Compact RAG box top-right: 4 dots labelled Budget / Planning / Resources / Overall."""
    panel_x, panel_y = Inches(7.0), Inches(0.95)
    panel_w, panel_h = Inches(6.0), Inches(1.6)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              panel_w - Inches(0.3), Inches(0.30),
              "STATUS", bold=True, size=11)

    rags = [
        ("Budget", data.rag_budget),
        ("Planning", data.rag_planning),
        ("Resources", data.rag_resources),
        ("Overall", data.rag_overall),
    ]
    col_w = (panel_w - Inches(0.3)) / 4
    for i, (label, status) in enumerate(rags):
        cx = panel_x + Inches(0.15) + col_w * i
        _add_rag_dot(slide, cx + col_w / 2 - Inches(0.125),
                     panel_y + Inches(0.45), status)
        _add_text(slide, cx, panel_y + Inches(0.85),
                  col_w, Inches(0.3),
                  label, bold=True, size=10, align=PP_ALIGN.CENTER)
        _add_text(slide, cx, panel_y + Inches(1.10),
                  col_w, Inches(0.3),
                  status.upper(), size=9,
                  align=PP_ALIGN.CENTER,
                  color=RAG_COLORS.get(status, TEXT_DARK))


def _build_timeline(slide, data: HighlightReportData) -> None:
    """Monthly timeline: 10 columns, one row per phase with a filled bar."""
    panel_x, panel_y = Inches(0.3), Inches(2.70)
    panel_w, panel_h = Inches(12.7), Inches(1.55)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              panel_w - Inches(0.3), Inches(0.3),
              "Timeline", bold=True, size=11)

    today = data.report_date or date.today()
    months = _months_window(today, n=10)
    grid_x = panel_x + Inches(2.0)
    grid_y = panel_y + Inches(0.45)
    grid_w = panel_w - Inches(2.1)
    grid_h = panel_h - Inches(0.55)
    col_w = grid_w / len(months)

    # Header row: month labels
    for i, m in enumerate(months):
        _add_text(slide, grid_x + col_w * i, grid_y - Inches(0.05),
                  col_w, Inches(0.25),
                  m.strftime("%b"), size=8, align=PP_ALIGN.CENTER,
                  color=TEXT_MUTED)

    phases = data.phases[:3] if data.phases else [
        {"name": "Prepare", "start": today, "end": today + timedelta(days=60)},
    ]
    row_h = (grid_h - Inches(0.25)) / max(len(phases), 1)
    for r, phase in enumerate(phases):
        # Row label
        _add_text(slide, panel_x + Inches(0.15),
                  grid_y + Inches(0.20) + row_h * r,
                  Inches(1.8), row_h,
                  phase.get("name", ""), size=10, bold=True)
        ps = phase.get("start")
        pe = phase.get("end")
        if not ps or not pe:
            continue
        # Map dates onto month columns
        def col_for(d: date) -> float:
            for idx, m in enumerate(months):
                next_m = months[idx + 1] if idx + 1 < len(months) else date(
                    m.year + (m.month // 12), (m.month % 12) + 1, 1
                )
                if m <= d < next_m:
                    frac = (d - m).days / max((next_m - m).days, 1)
                    return idx + frac
            if d < months[0]:
                return 0.0
            return float(len(months))
        c0 = max(0.0, col_for(ps))
        c1 = min(float(len(months)), col_for(pe))
        if c1 <= c0:
            continue
        bar_x = int(grid_x + c0 * col_w)
        bar_w = int((c1 - c0) * col_w)
        bar_y = int(grid_y + Inches(0.25) + row_h * r + Inches(0.05))
        bar_h = int(row_h - Inches(0.15))
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      bar_x, bar_y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = TIMELINE_FILL
        bar.line.fill.background()


def _build_financials(slide, data: HighlightReportData) -> None:
    panel_x, panel_y = Inches(0.3), Inches(4.40)
    panel_w, panel_h = Inches(6.2), Inches(2.4)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              panel_w - Inches(0.3), Inches(0.3),
              "Financials", bold=True, size=11)

    headers = ["", "Budget", "Actuals", "ETC", "Variance"]
    rows = data.financial_rows or [
        {"label": "Costs", "budget": 0, "actuals": 0, "etc": 0, "variance": 0},
        {"label": "Internal", "budget": 0, "actuals": 0, "etc": 0, "variance": 0},
        {"label": "External", "budget": 0, "actuals": 0, "etc": 0, "variance": 0},
        {"label": "Contingency", "budget": 0, "actuals": 0, "etc": 0, "variance": 0},
    ]
    cols = len(headers)
    col_w = (panel_w - Inches(0.3)) / cols
    top = panel_y + Inches(0.40)
    row_h = (panel_h - Inches(0.5)) / (len(rows) + 1)

    for ci, h in enumerate(headers):
        _add_text(slide, panel_x + Inches(0.15) + col_w * ci, top,
                  col_w, row_h, h, bold=True, size=9,
                  align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT,
                  color=TEXT_MUTED)
    for ri, row in enumerate(rows):
        y = top + row_h * (ri + 1)
        for ci, key in enumerate(["label", "budget", "actuals", "etc", "variance"]):
            v = row.get(key, "")
            if ci > 0 and isinstance(v, (int, float)):
                txt = f"€ {v:,.0f}"
            else:
                txt = str(v) if v is not None else ""
            _add_text(slide, panel_x + Inches(0.15) + col_w * ci, y,
                      col_w, row_h, txt, size=9,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT,
                      bold=(ci == 0))


def _build_risk_map(slide, data: HighlightReportData) -> None:
    """5x5 grid, probability (y) x impact (x). Cells coloured by risk density."""
    panel_x, panel_y = Inches(6.65), Inches(4.40)
    panel_w, panel_h = Inches(3.5), Inches(2.4)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              panel_w - Inches(0.3), Inches(0.3),
              "Risk Map (5x5)", bold=True, size=11)

    grid_x = panel_x + Inches(0.45)
    grid_y = panel_y + Inches(0.40)
    grid_w = panel_w - Inches(0.55)
    grid_h = panel_h - Inches(0.65)
    cell_w = grid_w / 5
    cell_h = grid_h / 5

    # Pre-bucket the risks into (impact, probability) -> count
    bucket = {}
    for r in data.risks_map:
        p = max(1, min(5, int(r.get("probability", 1))))
        i = max(1, min(5, int(r.get("impact", 1))))
        bucket[(i, p)] = bucket.get((i, p), 0) + 1

    for col in range(5):  # impact 1..5 along x
        for row in range(5):  # probability 1..5 along y (top = highest)
            impact = col + 1
            prob = 5 - row
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                grid_x + cell_w * col,
                grid_y + cell_h * row,
                cell_w, cell_h,
            )
            risk_score = impact * prob
            if risk_score >= 15:
                fill = RGBColor(0xEF, 0x9A, 0x9A)  # red zone
            elif risk_score >= 8:
                fill = RGBColor(0xFF, 0xE0, 0x82)  # amber zone
            else:
                fill = RGBColor(0xC8, 0xE6, 0xC9)  # green zone
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.line.width = Pt(1)
            count = bucket.get((impact, prob), 0)
            if count:
                tf = cell.text_frame
                tf.margin_left = tf.margin_right = Inches(0.02)
                p_ = tf.paragraphs[0]
                p_.alignment = PP_ALIGN.CENTER
                run = p_.add_run()
                run.text = str(count)
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = TEXT_DARK

    # Axis labels
    _add_text(slide, grid_x, grid_y + grid_h, grid_w, Inches(0.20),
              "Impact →", size=8, align=PP_ALIGN.CENTER, color=TEXT_MUTED)
    _add_text(slide, panel_x + Inches(0.05), grid_y, Inches(0.35), grid_h,
              "Probability →", size=8, color=TEXT_MUTED)


def _build_narrative(slide, data: HighlightReportData) -> None:
    panel_x, panel_y = Inches(10.30), Inches(4.40)
    panel_w, panel_h = Inches(2.7), Inches(2.4)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              panel_w - Inches(0.3), Inches(0.3),
              "Highlights / Lowlights", bold=True, size=11)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.35),
              panel_w - Inches(0.3), Inches(0.25),
              "Highlights", bold=True, size=9, color=RAG_COLORS["green"])
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.60),
              panel_w - Inches(0.3), Inches(0.70),
              data.highlights or "—", size=8)
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(1.35),
              panel_w - Inches(0.3), Inches(0.25),
              "Lowlights", bold=True, size=9, color=RAG_COLORS["red"])
    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(1.60),
              panel_w - Inches(0.3), Inches(0.70),
              data.lowlights or "—", size=8)


def _build_risks_issues(slide, data: HighlightReportData) -> None:
    panel_x, panel_y = Inches(0.3), Inches(6.90)
    panel_w, panel_h = Inches(12.7), Inches(0.55)
    _add_panel(slide, panel_x, panel_y, panel_w, panel_h)

    def join(items: Iterable[str], cap: int = 3) -> str:
        items = [s for s in items if s]
        if not items:
            return "—"
        head = items[:cap]
        extra = len(items) - cap
        text = " · ".join(head)
        if extra > 0:
            text += f" (+{extra} more)"
        return text

    _add_text(slide, panel_x + Inches(0.15), panel_y + Inches(0.05),
              Inches(1.0), Inches(0.40), "Issues:", bold=True, size=10)
    _add_text(slide, panel_x + Inches(1.15), panel_y + Inches(0.05),
              Inches(5.4), Inches(0.40), join(data.issues), size=10)
    _add_text(slide, panel_x + Inches(6.70), panel_y + Inches(0.05),
              Inches(1.0), Inches(0.40), "Risks:", bold=True, size=10)
    _add_text(slide, panel_x + Inches(7.70), panel_y + Inches(0.05),
              Inches(5.0), Inches(0.40), join(data.risks), size=10)


def _build_body(prs: Presentation, data: HighlightReportData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _build_header_block(slide, data)
    _build_rag_block(slide, data)
    _build_timeline(slide, data)
    _build_financials(slide, data)
    _build_risk_map(slide, data)
    _build_narrative(slide, data)
    _build_risks_issues(slide, data)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def render_highlight_pptx(data: HighlightReportData) -> bytes:
    """Render a 2-slide Highlight Report PPTX and return raw bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    _build_cover(prs, data)
    _build_body(prs, data)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Django adapter: HighlightReport ORM -> HighlightReportData
# ---------------------------------------------------------------------------

def highlight_report_to_data(report) -> HighlightReportData:
    """Convert a prince2.HighlightReport instance into the renderer DTO."""
    project = report.project

    # Resolve roles where we can; tolerate missing data.
    sponsor = ""
    senior_supplier = ""
    project_manager = ""
    try:
        board = getattr(project, "prince2_board", None)
        members = list(board.members.all()) if board else []
        for m in members:
            role = (m.role or "").lower()
            user = getattr(m, "user", None)
            name = (
                getattr(user, "first_name", "") or getattr(user, "email", "")
            ) if user else ""
            if role == "executive" and not sponsor:
                sponsor = name
            if role == "senior_supplier" and not senior_supplier:
                senior_supplier = name
            if role == "project_manager" and not project_manager:
                project_manager = name
    except Exception:
        pass

    # Phases from PRINCE2 Stage list (related_name = prince2_stages)
    phases = []
    try:
        for stage in project.prince2_stages.all().order_by("planned_start_date"):
            phases.append({
                "name": stage.name,
                "start": stage.planned_start_date,
                "end": stage.planned_end_date,
            })
    except Exception:
        pass

    # Financials from ProjectBudget
    fin_rows = []
    try:
        pb = project.project_budget
        total = float(pb.total_budget)
        spent = float(pb.total_spent)
        etc = float(pb.etc)
        cont = float(pb.contingency)
        fin_rows = [
            {"label": "Costs", "budget": total, "actuals": spent, "etc": etc,
             "variance": total - (spent + etc)},
            {"label": "Contingency", "budget": cont, "actuals": 0,
             "etc": 0, "variance": cont},
        ]
    except Exception:
        pass

    # Risk map
    risks_map = []
    try:
        from projects.models import Risk  # local import; avoid circular
        for r in Risk.objects.filter(project=project).only(
            "id", "name", "probability", "impact"
        ):
            # Project Risk.probability is 0-100; impact is High/Medium/Low.
            prob = max(1, min(5, round((r.probability or 0) / 20) or 1))
            impact_map = {"low": 1, "medium": 3, "high": 5}
            impact = impact_map.get((r.impact or "").lower(), 3)
            risks_map.append({"id": f"R{r.id}", "name": r.name,
                              "probability": prob, "impact": impact})
    except Exception:
        pass

    return HighlightReportData(
        project_name=project.name,
        report_date=report.report_date,
        period_label=(
            f"{report.period_start} → {report.period_end}"
            if report.period_start and report.period_end else ""
        ),
        report_month=report.report_date.month if report.report_date else None,
        report_to="Project Board",
        sponsor=sponsor,
        project_manager=project_manager,
        senior_supplier=senior_supplier,
        objectives=(project.project_goal or "")[:200],
        rag_overall=report.overall_status,
        rag_budget=report.rag_budget,
        rag_planning=report.rag_planning,
        rag_resources=report.rag_resources,
        phases=phases,
        financial_rows=fin_rows,
        risks_map=risks_map,
        highlights=report.highlights or report.work_completed or "",
        lowlights=report.lowlights or "",
        issues=[s for s in (report.issues_summary or "").splitlines() if s.strip()],
        risks=[s for s in (report.risks_summary or "").splitlines() if s.strip()],
    )
